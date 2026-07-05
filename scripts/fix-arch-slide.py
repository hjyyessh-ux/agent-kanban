#!/usr/bin/env python3
"""Replace the architecture *image* on the '설계 ①' slide with native
PowerPoint shapes (editable vector boxes/text), preserving every other
slide and the user's manual edits.

Usage: python scripts/fix-arch-slide.py [path.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TARGET = sys.argv[1] if len(sys.argv) > 1 else "outputs/agent-kanban-deck-origin.pptx"

# ── palette (matches the deck) ───────────────────────────────────────
BG      = RGBColor(0x0F, 0x12, 0x18)
PANEL   = RGBColor(0x17, 0x1C, 0x26)
INK     = RGBColor(0xE7, 0xEC, 0xF3)
MUT     = RGBColor(0x9A, 0xA6, 0xB8)
BLUE    = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
PINK    = RGBColor(0xF4, 0x72, 0xB6)
CARD_BD = RGBColor(0x2B, 0x33, 0x42)
MONO_INK = RGBColor(0xCD, 0xD6, 0xE4)
FONT = "Apple SD Gothic Neo"
MONO = "Menlo"


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         font=FONT, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for t, c, b in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b
        r.font.color.rgb = c; r.font.name = font
    return p


def _rbox(s, x, y, w, h, accent=CARD_BD, fill=PANEL, round=0.06):
    b = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        b.adjustments[0] = round
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = accent; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    return b


def _txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    return textbox(s, Inches(x), Inches(y), Inches(w), Inches(h), anchor)


def _arrow(s, x, y, w, ch, color=BLUE, size=26, align=PP_ALIGN.CENTER):
    tf = _txt(s, x, y, w, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, ch, size, color, bold=True, first=True, align=align)


def _chip(s, x, y, w, h, text, fontcolor=INK, border=CARD_BD, mono=False,
          size=10.5, fill=BG, bold=False, dot=None):
    b = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        b.adjustments[0] = 0.28
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = border; b.line.width = Pt(1)
    b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.09); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    if dot:
        rd = p.add_run(); rd.text = "● "; rd.font.size = Pt(size)
        rd.font.color.rgb = dot; rd.font.name = FONT
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = fontcolor; r.font.name = (MONO if mono else FONT)
    return b


def draw_arch(s):
    # ── Row 1 ──────────────────────────────────────────────────────
    _rbox(s, 0.55, 1.72, 3.25, 1.98)
    tf = _txt(s, 0.76, 1.86, 2.9, 0.72)
    para(tf, "① Agent 런타임 (3종)", 14, INK, bold=True, first=True, space_after=2)
    para(tf, "각자 잘하는 모델을 카드 단위로 선택", 10.5, MUT, space_after=0)
    _chip(s, 0.78, 2.64, 2.82, 0.33, "opencode · session.prompt()", INK, GREEN, dot=GREEN, size=10.5)
    _chip(s, 0.78, 3.02, 2.82, 0.33, "codex · codex CLI", INK, PURPLE, dot=PURPLE, size=10.5)
    _chip(s, 0.78, 3.40, 2.82, 0.33, "claude · claude CLI", INK, AMBER, dot=AMBER, size=10.5)

    _arrow(s, 3.80, 2.42, 0.55, "→", BLUE, 26)

    _rbox(s, 4.40, 1.72, 8.35, 1.98, accent=BLUE)
    tf = _txt(s, 4.62, 1.86, 8.0, 0.72)
    para(tf, [("② Plugin 런타임    ", INK, True), ("@opencode-ai/plugin", BLUE, False)],
         14, INK, first=True, space_after=2)
    para(tf, "프롬프트 디스패치 · 훅으로 응답 수집 · 백그라운드 모니터", 10.5, MUT, space_after=0)
    subs = [("Runtime Registry", "runtime별 adapter 선택\n세션 생성 / 재사용"),
            ("Event Hooks", "response / idle 훅\n카드 자동 완료 · 큐 디스패치"),
            ("Monitors", "Scheduler · Telegram\nStale / Question 감시")]
    for i, (t, d) in enumerate(subs):
        x = 4.62 + i * 2.72
        _rbox(s, x, 2.64, 2.56, 0.98)
        tf = _txt(s, x + 0.14, 2.72, 2.3, 0.85)
        para(tf, t, 11.5, INK, bold=True, first=True, space_after=2)
        para(tf, d, 9.5, MUT, space_after=0)

    # ── middle band ────────────────────────────────────────────────
    _arrow(s, 0.55, 3.80, 12.2, "↓   Bun.serve() HTTP · :24680   (REST + 정적 SPA)", BLUE, 15)

    # ── Row 2 ──────────────────────────────────────────────────────
    _rbox(s, 0.55, 4.38, 6.05, 2.12, accent=PINK)
    tf = _txt(s, 0.78, 4.52, 5.7, 0.7)
    para(tf, [("③ 로컬 저장소    ", INK, True), ("~/.agent-kanban", PINK, False)],
         14, INK, first=True, space_after=2)
    para(tf, "도메인별 JSON · 원자적 쓰기 · 듀얼 락", 10.5, MUT, space_after=0)
    files = ["active.json", "archive/YYYY-MM.json", "schedulers.json",
             "scripts.json", "settings.json", "screenshots/"]
    for i, fn in enumerate(files):
        col, row = i % 3, i // 3
        _chip(s, 0.78 + col * 1.92, 5.34 + row * 0.46, 1.82, 0.34, fn,
              MONO_INK, CARD_BD, mono=True, size=9.5)

    _arrow(s, 6.62, 5.12, 0.5, "↔", BLUE, 24)

    _rbox(s, 7.15, 4.38, 5.6, 2.12, accent=GREEN)
    tf = _txt(s, 7.37, 4.52, 5.2, 0.7)
    para(tf, "④ React SPA (Vite)", 14, INK, bold=True, first=True, space_after=2)
    para(tf, "3초 폴링 동기화 · 보드 / 스케줄러 / 스크립트 / 설정 탭", 10.5, MUT, space_after=0)
    tags = ["Board · Session 모아보기", "Card Detail · Queue · Resume", "Screenshot 첨부"]
    for i, tg in enumerate(tags):
        _chip(s, 7.37, 5.34 + i * 0.40, 5.0, 0.34, tg, INK, CARD_BD, size=10.5)

    tf = _txt(s, 0.55, 6.66, 12.2, 0.5)
    para(tf, [("●  카드 = 작업 단위 (제목·설명·런타임·모델·세션)        ", MUT, False),
              ("●  모든 데이터는 로컬 파일 — 외부 의존성 없음", MUT, False)],
         11, MUT, first=True)


def main():
    prs = Presentation(TARGET)
    target = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and "아키텍처" in sh.text_frame.text:
                target = s
                break
        if target is not None:
            break
    if target is None:
        raise SystemExit("아키텍처 슬라이드를 찾지 못했습니다.")

    removed = 0
    for sh in list(target.shapes):
        if sh.shape_type == 13:  # PICTURE
            sh._element.getparent().remove(sh._element)
            removed += 1

    draw_arch(target)
    prs.save(TARGET)
    print(f"saved: {TARGET}  (removed {removed} picture(s), drew native architecture)")


if __name__ == "__main__":
    main()
