from pptx import Presentation
import sys
f = sys.argv[1]
p = Presentation(f)
def inch(e):
    try: return round(e/914400,2)
    except: return e
print(f"FILE={f}  slides={len(list(p.slides))}")
for i,s in enumerate(p.slides,1):
    print(f"\n=== slide {i} ===")
    for sh in s.shapes:
        try: geo=f"L{inch(sh.left)} T{inch(sh.top)} W{inch(sh.width)} H{inch(sh.height)}"
        except: geo="?"
        txt=""
        if sh.has_text_frame:
            txt = sh.text_frame.text.replace("\n"," / ")[:80]
        print(f"  [{sh.shape_type}] {geo}  {txt}")
