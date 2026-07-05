#!/usr/bin/env python3
"""Insert a brief 'code architecture / directory layout' slide as the 4th
slide (before the system-architecture slide). Idempotent: skips if already
present.

Usage: python scripts/add-code-arch-slide.py [path.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TARGET = sys.argv[1] if len(sys.argv) > 1 else "outputs/agent-kanban-deck-origin.pptx"
INSERT_AT = 3  # 0-based -> becomes the 4th slide
MARKER = "코드 구조"

BG      = RGBColor(0x0F, 0x12, 0x18)
PANEL   = RGBColor(0x17, 0x1C, 0x26)
INK     = RGBColor(0xE7, 0xEC, 0xF3)
MUT     = RGBColor(0x9A, 0xA6, 0xB8)
BLUE    = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
PINK    = RGBColor(0xF4, 0x72, 0xB6)
CARD_BD = RGBColor(0x2B, 0x33, 0x42)
FONT = "Apple SD Gothic Neo"
MONO = "Menlo"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         font=FONT, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for t, c, b in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b
        r.font.color.rgb = c; r.font.name = font
    return p


def accent_bar(s, color=PURPLE):
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.16), EMU_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False


def kicker_title(s, kicker, title, color=PURPLE):
    accent_bar(s, color)
    tf = textbox(s, Inches(0.62), Inches(0.42), Inches(12.1), Inches(1.1))
    para(tf, kicker.upper(), 13, color, bold=True, first=True, space_after=2)
    para(tf, title, 30, INK, bold=True, space_after=0)


def layer_row(s, y, dirname, accent, desc):
    b = s.shapes.add_shape(5, Inches(0.62), Inches(y), Inches(12.1), Inches(0.86))
    try:
        b.adjustments[0] = 0.12
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = PANEL
    b.line.color.rgb = accent; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    tf = textbox(s, Inches(0.95), Inches(y), Inches(11.5), Inches(0.86), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = dirname
    r.font.size = Pt(15); r.font.bold = True; r.font.name = MONO; r.font.color.rgb = accent
    r2 = p.add_run(); r2.text = "    " + desc
    r2.font.size = Pt(12.5); r2.font.color.rgb = INK; r2.font.name = FONT


def build_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)

    kicker_title(s, "코드 구조", "코드는 이렇게 나뉘어 있다", PURPLE)

    rows = [
        ("src/core/",  PINK,  "공유 타입 · 스토어 · 파일 락 · 에이전트 설정  —  types.ts 가 단일 소스"),
        ("src/plugin/", BLUE, "툴 · 훅 · 디스패치 · 런타임 어댑터(claude·codex·opencode) · Telegram · 스케줄러"),
        ("src/server/", AMBER, "Bun.serve() HTTP · REST 라우트 · 정적 SPA 서빙"),
        ("web/src/",   GREEN, "React SPA · App.tsx 탭/모달 · 훅이 fetch/폴링 · 플레인 CSS"),
    ]
    y = 1.95
    for dirname, accent, desc in rows:
        layer_row(s, y, dirname, accent, desc)
        y += 1.02

    # aux line
    tf = textbox(s, Inches(0.95), Inches(y + 0.02), Inches(11.5), Inches(0.4))
    para(tf, [("보조    ", MUT, True),
              ("e2e/ — Playwright E2E      ·      scripts/ — 운영 CLI 스크립트", MUT, False)],
         12, MUT, first=True)

    # stack line
    tf = textbox(s, Inches(0.95), Inches(y + 0.55), Inches(11.5), Inches(0.4))
    para(tf, [("STACK    ", PURPLE, True),
              ("Bun · TypeScript · React 19 · Vite · zod · 로컬 JSON 파일 영속화", INK, False)],
         13, INK, first=True)
    return s


def main():
    prs = Presentation(TARGET)
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and MARKER in sh.text_frame.text:
                print("이미 코드 구조 슬라이드가 있어 건너뜁니다.")
                return

    build_slide(prs)

    # move the just-added (last) slide to INSERT_AT
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    sld_lst.remove(ids[-1])
    sld_lst.insert(INSERT_AT, ids[-1])

    prs.save(TARGET)
    print(f"saved: {TARGET}  (inserted '코드 구조' slide at position {INSERT_AT + 1})")


if __name__ == "__main__":
    main()
