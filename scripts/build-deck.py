#!/usr/bin/env python3
"""Build the Agent Kanban presentation deck (outputs/agent-kanban-deck.pptx).

Dark, screenshot-first deck. Minimal on-slide text — the presenter narrates.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(HERE, "outputs", "slides", "img")
OUT = os.path.join(HERE, "outputs", "agent-kanban-deck.pptx")

# ── palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x12, 0x18)
PANEL    = RGBColor(0x17, 0x1C, 0x26)
INK      = RGBColor(0xE7, 0xEC, 0xF3)
MUT      = RGBColor(0x9A, 0xA6, 0xB8)
BLUE     = RGBColor(0x38, 0xBD, 0xF8)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE   = RGBColor(0xA7, 0x8B, 0xFA)
PINK     = RGBColor(0xF4, 0x72, 0xB6)
CARD_BD  = RGBColor(0x2B, 0x33, 0x42)
FONT = "Apple SD Gothic Neo"
MONO = "Menlo"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         font=FONT, space_after=6, bullet=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for i, item in enumerate(runs):
        t, c, b = item
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b
        r.font.color.rgb = c; r.font.name = font
    return p


def chip(s, x, y, text, color, w=Inches(2.0), h=Inches(0.42)):
    box = s.shapes.add_shape(5, x, y, w, h)  # rounded rect
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = color; box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = INK; r.font.name = FONT
    return box


def panel(s, x, y, w, h, fill=PANEL, line=CARD_BD):
    box = s.shapes.add_shape(5, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = line; box.line.width = Pt(1)
    box.shadow.inherit = False
    return box


def accent_bar(s, color=BLUE):
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.16), EMU_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False


def kicker_title(s, kicker, title, color=BLUE):
    accent_bar(s, color)
    tf = textbox(s, Inches(0.62), Inches(0.42), Inches(12.1), Inches(1.1))
    para(tf, kicker.upper(), 13, color, bold=True, first=True, space_after=2)
    para(tf, title, 30, INK, bold=True, space_after=0)


def add_image(s, img_path, x, y, max_w, max_h, frame=True):
    """Place image fitted into the (max_w, max_h) box, centered, with a border."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    px = int(x + (max_w - w) / 2)
    py = int(y + (max_h - h) / 2)
    if frame:
        bd = panel(s, Emu(px - Emu(0).emu), Emu(py), Emu(w), Emu(h), fill=PANEL, line=CARD_BD)
    pic = s.shapes.add_picture(img_path, Emu(px), Emu(py), width=Emu(w), height=Emu(h))
    pic.line.color.rgb = CARD_BD; pic.line.width = Pt(1)
    return px, py, w, h


def caption(s, x, y, w, text, color=MUT, size=12, align=PP_ALIGN.CENTER):
    tf = textbox(s, x, y, w, Inches(0.5))
    para(tf, text, size, color, first=True, align=align)


def P(img):
    return os.path.join(IMG, img)


# ── native architecture diagram (vector shapes, not a screenshot) ─────
MONO_INK = RGBColor(0xCD, 0xD6, 0xE4)


def _rbox(s, x, y, w, h, accent=CARD_BD, fill=PANEL, round=0.06):
    b = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        b.adjustments[0] = round
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = accent; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    return b


def _txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    return textbox(s, Inches(x), Inches(y), Inches(w), Inches(h), anchor)


def _arrow(s, x, y, w, ch, color=BLUE, size=26, align=PP_ALIGN.CENTER):
    tf = _txt(s, x, y, w, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, ch, size, color, bold=True, first=True, align=align)


def _chip(s, x, y, w, h, text, fontcolor=INK, border=CARD_BD, mono=False,
          size=10.5, fill=BG, bold=False, dot=None):
    b = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        b.adjustments[0] = 0.28
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = border; b.line.width = Pt(1)
    b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.09); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    if dot:
        rd = p.add_run(); rd.text = "● "; rd.font.size = Pt(size)
        rd.font.color.rgb = dot; rd.font.name = FONT
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = fontcolor; r.font.name = (MONO if mono else FONT)
    return b


def draw_arch(s):
    # ── Row 1 ──────────────────────────────────────────────────────
    # ① Agent runtime
    _rbox(s, 0.55, 1.72, 3.25, 1.98)
    tf = _txt(s, 0.76, 1.86, 2.9, 0.72)
    para(tf, "① Agent 런타임 (3종)", 14, INK, bold=True, first=True, space_after=2)
    para(tf, "각자 잘하는 모델을 카드 단위로 선택", 10.5, MUT, space_after=0)
    _chip(s, 0.78, 2.64, 2.82, 0.33, "opencode · session.prompt()", INK, GREEN, dot=GREEN, size=10.5)
    _chip(s, 0.78, 3.02, 2.82, 0.33, "codex · codex CLI", INK, PURPLE, dot=PURPLE, size=10.5)
    _chip(s, 0.78, 3.40, 2.82, 0.33, "claude · claude CLI", INK, AMBER, dot=AMBER, size=10.5)

    _arrow(s, 3.80, 2.42, 0.55, "→", BLUE, 26)

    # ② Plugin runtime
    _rbox(s, 4.40, 1.72, 8.35, 1.98, accent=BLUE)
    tf = _txt(s, 4.62, 1.86, 8.0, 0.72)
    para(tf, [("② Plugin 런타임    ", INK, True), ("@opencode-ai/plugin", BLUE, False)],
         14, INK, first=True, space_after=2)
    para(tf, "프롬프트 디스패치 · 훅으로 응답 수집 · 백그라운드 모니터", 10.5, MUT, space_after=0)
    subs = [("Runtime Registry", "runtime별 adapter 선택\n세션 생성 / 재사용"),
            ("Event Hooks", "response / idle 훅\n카드 자동 완료 · 큐 디스패치"),
            ("Monitors", "Scheduler · Telegram\nStale / Question 감시")]
    for i, (t, d) in enumerate(subs):
        x = 4.62 + i * 2.72
        _rbox(s, x, 2.64, 2.56, 0.98)
        tf = _txt(s, x + 0.14, 2.72, 2.3, 0.85)
        para(tf, t, 11.5, INK, bold=True, first=True, space_after=2)
        para(tf, d, 9.5, MUT, space_after=0)

    # ── middle band ────────────────────────────────────────────────
    _arrow(s, 0.55, 3.80, 12.2, "↓   Bun.serve() HTTP · :24680   (REST + 정적 SPA)", BLUE, 15)

    # ── Row 2 ──────────────────────────────────────────────────────
    # ③ storage
    _rbox(s, 0.55, 4.38, 6.05, 2.12, accent=PINK)
    tf = _txt(s, 0.78, 4.52, 5.7, 0.7)
    para(tf, [("③ 로컬 저장소    ", INK, True), ("~/.agent-kanban", PINK, False)],
         14, INK, first=True, space_after=2)
    para(tf, "도메인별 JSON · 원자적 쓰기 · 듀얼 락", 10.5, MUT, space_after=0)
    files = ["active.json", "archive/YYYY-MM.json", "schedulers.json",
             "scripts.json", "settings.json", "screenshots/"]
    for i, fn in enumerate(files):
        col, row = i % 3, i // 3
        _chip(s, 0.78 + col * 1.92, 5.34 + row * 0.46, 1.82, 0.34, fn,
              MONO_INK, CARD_BD, mono=True, size=9.5)

    _arrow(s, 6.62, 5.12, 0.5, "↔", BLUE, 24)

    # ④ React SPA
    _rbox(s, 7.15, 4.38, 5.6, 2.12, accent=GREEN)
    tf = _txt(s, 7.37, 4.52, 5.2, 0.7)
    para(tf, "④ React SPA (Vite)", 14, INK, bold=True, first=True, space_after=2)
    para(tf, "3초 폴링 동기화 · 보드 / 스케줄러 / 스크립트 / 설정 탭", 10.5, MUT, space_after=0)
    tags = ["Board · Session 모아보기", "Card Detail · Queue · Resume", "Screenshot 첨부"]
    for i, tg in enumerate(tags):
        _chip(s, 7.37, 5.34 + i * 0.40, 5.0, 0.34, tg, INK, CARD_BD, size=10.5)

    # legend
    tf = _txt(s, 0.55, 6.66, 12.2, 0.5)
    para(tf, [("●  카드 = 작업 단위 (제목·설명·런타임·모델·세션)        ", MUT, False),
              ("●  모든 데이터는 로컬 파일 — 외부 의존성 없음", MUT, False)],
         11, MUT, first=True)


# ════════════════════════════════════════════════════════════════════
# 1. TITLE
# ════════════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, BLUE)
tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
para(tf, "Agent Kanban", 54, INK, bold=True, first=True, space_after=8)
para(tf, [("여러 AI 에이전트를 ", INK, False), ("한 보드에서", BLUE, True), (" 관리하기", INK, False)],
     26, INK, space_after=18)
para(tf, "opencode · codex · claude — 하나의 칸반으로", 18, MUT, space_after=0)
# runtime chips
chip(s, Inches(0.92), Inches(5.0), "●  opencode", GREEN, w=Inches(2.1))
chip(s, Inches(3.12), Inches(5.0), "●  codex", PURPLE, w=Inches(1.9))
chip(s, Inches(5.12), Inches(5.0), "●  claude", AMBER, w=Inches(1.9))
tf = textbox(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
para(tf, "사내 발표 · 오픈소스 기여 제안", 13, MUT, first=True)

# ════════════════════════════════════════════════════════════════════
# 2. WHY — motivation (text-forward)
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "왜 만들었나", "한 번에 한 가지 일만 하지 않는 시대", AMBER)
# left: pain  right: turn
panel(s, Inches(0.62), Inches(1.95), Inches(5.95), Inches(4.9), fill=PANEL)
tf = textbox(s, Inches(0.92), Inches(2.2), Inches(5.4), Inches(4.4))
para(tf, "불편했던 점", 17, AMBER, bold=True, first=True, space_after=10)
for t in [
    "TUI·CLI 여러 개를 동시에 띄워 작업",
    "어느 순간 헷갈리고 잊어버린다",
    "CLI는 이력 관리가 안 되고, Desktop은 남아도 보기 불편",
    "tmux / iTerm2 단축키를 다 외워야 했다",
]:
    para(tf, [("•  ", MUT, False), (t, INK, False)], 15, INK, space_after=9)

panel(s, Inches(6.78), Inches(1.95), Inches(5.95), Inches(4.9), fill=PANEL, line=BLUE)
tf = textbox(s, Inches(7.08), Inches(2.2), Inches(5.4), Inches(4.4))
para(tf, "그래서 — 핵심 한 줄", 17, BLUE, bold=True, first=True, space_after=10)
para(tf, [("여러 세션이 있어도 ", INK, False), ("결국 한 화면에서", BLUE, True),
          (" 본다", INK, False)], 19, INK, bold=True, space_after=14)
for t in [
    "메모앱처럼 생각날 때 카드로 기록 → 나중에 실행",
    "컨텍스트를 옮길 필요 없이 그 자리에서 TODO",
    "내 취향대로 기능을 직접 붙일 수 있다",
]:
    para(tf, [("•  ", MUT, False), (t, INK, False)], 15, INK, space_after=9)

# ════════════════════════════════════════════════════════════════════
# 3. THE ANSWER — board full
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "결과물", "여러 에이전트의 작업을 한 보드에서", BLUE)
add_image(s, P("01-board-full.png"), Inches(0.62), Inches(1.7), Inches(12.1), Inches(5.0))
caption(s, Inches(0.62), Inches(6.85), Inches(12.1),
        "To Do · In Progress · Complete · Done — 런타임/모델/세션이 카드마다 표시")

# ════════════════════════════════════════════════════════════════════
# 4. DESIGN — architecture
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "설계 ①", "아키텍처 — 3 런타임 · 1 플러그인 · 로컬 JSON", PURPLE)
draw_arch(s)

# ════════════════════════════════════════════════════════════════════
# 5. DESIGN — execution flow
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "설계 ②", "카드 한 장이 실행되는 과정", PURPLE)
add_image(s, P("00b-flow.png"), Inches(0.62), Inches(1.75), Inches(12.1), Inches(5.4), frame=False)

# ════════════════════════════════════════════════════════════════════
# 6. FEATURE — session 모아보기 + workspace
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "기능", "세션 모아보기 · 워크스페이스 모아보기", GREEN)
add_image(s, P("02-session-group.png"), Inches(0.62), Inches(1.8), Inches(8.4), Inches(4.7))
caption(s, Inches(0.62), Inches(6.6), Inches(8.4),
        "완료 카드를 세션 단위로 묶어 대화 맥락 그대로 다시 보기", align=PP_ALIGN.CENTER)
panel(s, Inches(9.25), Inches(1.95), Inches(3.5), Inches(2.2), fill=PANEL)
tf = textbox(s, Inches(9.5), Inches(2.15), Inches(3.0), Inches(1.9))
para(tf, "워크스페이스 모아보기", 14, GREEN, bold=True, first=True, space_after=8)
para(tf, "디렉토리별로 카드를 필터링 — 프로젝트가 섞이지 않는다", 13, MUT, space_after=0)
add_image(s, P("03-workspace-switcher.png"), Inches(9.25), Inches(4.35), Inches(3.5), Inches(0.7), frame=False)

# ════════════════════════════════════════════════════════════════════
# 7. FEATURE — Session Resume
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "기능 · 이어하기", "Session Resume — 그 대화에서 이어서", AMBER)
add_image(s, P("07-resume-panel.png"), Inches(7.0), Inches(1.65), Inches(5.9), Inches(5.0))
tf = textbox(s, Inches(0.62), Inches(2.0), Inches(6.1), Inches(5.0))
para(tf, "기존 세션을 골라 그 컨텍스트 그대로 새 작업을 시작", 17, INK, bold=True, first=True, space_after=16)
para(tf, "예시", 14, AMBER, bold=True, space_after=6)
for t in [
    "어제 “결제 모듈 리팩터링”을 claude로 진행 → 완료",
    "오늘 “방금 그 리팩터링에 테스트 추가” 카드 생성",
    "Resume로 어제 세션 선택 → 코드 맥락 다시 설명 불필요",
]:
    para(tf, [("•  ", MUT, False), (t, INK, False)], 14, INK, space_after=8)
para(tf, "실패한 작업도 같은 세션에서 자동 재시도", 13, MUT, space_after=0)
caption(s, Inches(7.0), Inches(6.8), Inches(5.9), "todo 카드 상세 → Session Resume 목록에서 SELECT")

# ════════════════════════════════════════════════════════════════════
# 8. FEATURE — Queue After
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "기능 · 큐", "Queue After — 끝나면 다음 작업 자동 실행", BLUE)
add_image(s, P("06-queue-panel.png"), Inches(0.62), Inches(1.85), Inches(6.5), Inches(4.4))
add_image(s, P("08-queued-card.png"), Inches(10.7), Inches(3.6), Inches(2.1), Inches(2.6))
tf = textbox(s, Inches(7.3), Inches(1.95), Inches(3.2), Inches(4.4))
para(tf, "앞 카드가 끝나면 다음 카드를 자동 디스패치", 16, INK, bold=True, first=True, space_after=14)
para(tf, "두 가지 모드", 14, BLUE, bold=True, space_after=6)
para(tf, [("• New Session\n", INK, True), ("  깨끗한 새 세션으로 시작", MUT, False)], 13, INK, space_after=8)
para(tf, [("• Continue After\n", INK, True), ("  앞 카드 세션을 이어받음", MUT, False)], 13, INK, space_after=12)
para(tf, "예시", 14, BLUE, bold=True, space_after=4)
para(tf, "“스키마 마이그레이션” 완료 직후 “마이그레이션 검증” 이 자동 실행 — 밤사이 체이닝", 13, MUT, space_after=0)
caption(s, Inches(0.62), Inches(6.35), Inches(6.5), "Target 선택 → 큐에 연결되면 보드에 ⏭ NEXT 배지")

# ════════════════════════════════════════════════════════════════════
# 9. FEATURE — Card detail / create
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "기능", "카드 = 작업 단위 (생성 · 상세)", GREEN)
add_image(s, P("05-create-dialog.png"), Inches(0.62), Inches(1.85), Inches(5.95), Inches(4.6))
add_image(s, P("04-card-detail.png"), Inches(6.78), Inches(1.85), Inches(5.95), Inches(4.6))
caption(s, Inches(0.62), Inches(6.6), Inches(5.95), "생성: 런타임·모델·디렉토리·스크린샷 첨부")
caption(s, Inches(6.78), Inches(6.6), Inches(5.95), "상세: 프롬프트 · 결과 · 세션 정보 한눈에")

# ════════════════════════════════════════════════════════════════════
# 10. FEATURE — Scheduler + Scripts
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "기능 · 자동화", "Scheduler · Script — 반복 작업을 자산으로", PINK)
add_image(s, P("09-scheduler.png"), Inches(0.62), Inches(1.85), Inches(5.95), Inches(4.5))
add_image(s, P("10-scripts.png"), Inches(6.78), Inches(1.85), Inches(5.95), Inches(4.5))
caption(s, Inches(0.62), Inches(6.5), Inches(5.95), "Scheduler: cron으로 정기 작업 (예: 평일 9시 PR 다이제스트)")
caption(s, Inches(6.78), Inches(6.5), Inches(5.95), "Script: 자주 쓰는 작업을 등록·실행, 히스토리 보관")

# ════════════════════════════════════════════════════════════════════
# 11. INTEGRATIONS — placeholders (user fills)
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "연동", "외부에서도 카드로 — Telegram · Google Task", BLUE)
tf = textbox(s, Inches(0.62), Inches(1.7), Inches(12.1), Inches(0.6))
para(tf, "갑자기 떠오른 아이디어/이슈를 외부 채널에서 바로 카드로 등록하고, 모바일에서 작업까지", 15, MUT, first=True)
# two placeholder frames
for i, (label, color) in enumerate([("Telegram", BLUE), ("Google Task", GREEN)]):
    x = Inches(0.62 + i * 6.16)
    box = panel(s, x, Inches(2.5), Inches(5.95), Inches(4.0), fill=PANEL, line=color)
    # dashed look via inner note
    tf = textbox(s, x, Inches(2.5), Inches(5.95), Inches(4.0), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, 22, color, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=10)
    para(tf, "📷 캡쳐 자리 — 직접 추가", 15, MUT, align=PP_ALIGN.CENTER, space_after=6)
    note = ("• 텔레그램에서 Model/Directory/Runtime 지정 후\n  Claude Code처럼 대화형 작업"
            if i == 0 else
            "• “KANBAN에 ~~ 작업 카드로 등록해줘”\n  → Gemini Task 연동으로 카드 자동 생성")
    para(tf, note, 13, MUT, align=PP_ALIGN.CENTER, space_after=0)

# ════════════════════════════════════════════════════════════════════
# 12. LIMITATIONS
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "솔직한 한계", "아직 부족한 점", PINK)
items = [
    ("Desktop 앱과 미동기화", "카드는 CLI/플러그인 방식으로 실행 → Claude·Codex Desktop 세션 이력과는 동기화되지 않음"),
    ("opencode 우선 설계", "처음부터 opencode 기준으로 시작 → codex·claude 지원은 기본 동작은 안정적이나 일부 기능은 미흡"),
    ("로컬 단일 사용자", "데이터가 로컬 JSON에 저장 → 팀 공유/멀티 디바이스 동기화는 아직 없음"),
    ("실행 환경 의존", "각 런타임 CLI가 로컬에 설치·로그인되어 있어야 동작"),
]
y = 1.95
for title, desc in items:
    panel(s, Inches(0.62), Inches(y), Inches(12.1), Inches(1.12), fill=PANEL)
    tf = textbox(s, Inches(0.95), Inches(y + 0.14), Inches(11.5), Inches(0.9))
    para(tf, title, 16, PINK, bold=True, first=True, space_after=3)
    para(tf, desc, 13, MUT, space_after=0)
    y += 1.28

# ════════════════════════════════════════════════════════════════════
# 13. FUTURE — LLM Wiki
# ════════════════════════════════════════════════════════════════════
s = slide()
kicker_title(s, "나아갈 방향", "다음은 LLM Wiki", GREEN)
panel(s, Inches(0.62), Inches(1.95), Inches(12.1), Inches(2.0), fill=PANEL, line=GREEN)
tf = textbox(s, Inches(0.95), Inches(2.2), Inches(11.5), Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("핵심 아이디어 — ", INK, True),
          ("Done에서 Archive해도 파일을 지우지 않는다", GREEN, True)], 19, INK, first=True, space_after=8)
para(tf, "그 이력을 WIKI로 남겨 재활용하면? (Andrej Karpathy의 LLM Wiki 컨셉)", 15, MUT, space_after=0)
# next steps row
nexts = [
    ("작업 이력 → Wiki", "지금까지의 작업을 Obsidian에 견고하게 기록"),
    ("멀티 에이전트 확장", "여러 에이전트를 spawn해 병렬 작업 강화"),
    ("Script/Scheduler 고도화", "자주 쓰는 패턴을 참조해 자동 수행"),
]
for i, (t, d) in enumerate(nexts):
    x = Inches(0.62 + i * 4.1)
    panel(s, x, Inches(4.25), Inches(3.9), Inches(2.4), fill=PANEL)
    tf = textbox(s, x + Inches(0.25), Inches(4.5), Inches(3.4), Inches(2.0))
    para(tf, f"0{i+1}", 22, GREEN, bold=True, first=True, space_after=6)
    para(tf, t, 15, INK, bold=True, space_after=6)
    para(tf, d, 13, MUT, space_after=0)

# ════════════════════════════════════════════════════════════════════
# 14. CLOSING
# ════════════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, BLUE)
tf = textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.5))
para(tf, "한곳에서 관리하면, 확장의 여지가 생긴다", 34, INK, bold=True, first=True, space_after=16)
para(tf, [("여러 에이전트의 결과가 한 보드에 모이니 ", MUT, False),
          ("자동화·Wiki·멀티에이전트", BLUE, True),
          ("로 뻗어나갈 수 있습니다.", MUT, False)], 18, MUT, space_after=0)
tf = textbox(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.0))
para(tf, "같이 써보고, 같이 만들어요 — 오픈소스 기여 환영", 20, INK, bold=True, first=True, space_after=6)
para(tf, "github · agent-kanban", 14, MUT, space_after=0)

prs.save(OUT)
print(f"saved: {OUT}  ({len(prs.slides.__iter__().__length_hint__() if False else list(prs.slides))} slides)")
