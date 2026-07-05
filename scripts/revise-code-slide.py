#!/usr/bin/env python3
"""Rewrite the 4th slide into a brief 'storage + stack' page:
where files live, how they are persisted, and the full tech stack.
No per-directory breakdown.

Usage: python scripts/revise-code-slide.py [path.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TARGET = sys.argv[1] if len(sys.argv) > 1 else "outputs/agent-kanban-deck-origin.pptx"
SLIDE_INDEX = 3  # 0-based -> the 4th slide

BG      = RGBColor(0x0F, 0x12, 0x18)
PANEL   = RGBColor(0x17, 0x1C, 0x26)
INK     = RGBColor(0xE7, 0xEC, 0xF3)
MUT     = RGBColor(0x9A, 0xA6, 0xB8)
BLUE    = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
PINK    = RGBColor(0xF4, 0x72, 0xB6)
CARD_BD = RGBColor(0x2B, 0x33, 0x42)
MONO_INK = RGBColor(0xCD, 0xD6, 0xE4)
FONT = "Apple SD Gothic Neo"
MONO = "Menlo"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def add_line(tf, parts, first=False, align=PP_ALIGN.LEFT, space_after=4, bullet_color=None):
    """parts = list of (text, size, color, bold, font)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after); p.space_before = Pt(0)
    for t, size, color, bold, font in parts:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return p


def panel(s, x, y, w, h, accent=CARD_BD, round=0.05):
    b = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        b.adjustments[0] = round
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = PANEL
    b.line.color.rgb = accent; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    return b


def rebuild(s):
    # clear everything
    spTree = s.shapes._spTree
    for sh in list(s.shapes):
        spTree.remove(sh._element)

    # background + accent bar
    bg = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    spTree.remove(bg._element); spTree.insert(2, bg._element)
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.16), EMU_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background(); bar.shadow.inherit = False

    # title
    tf = textbox(s, 0.62, 0.42, 12.1, 1.1)
    add_line(tf, [("코드 구조 · 데이터", 13, PURPLE, True, FONT)], first=True, space_after=2)
    add_line(tf, [("파일이 어디에 · 어떻게 저장되나", 30, INK, True, FONT)])

    # ── left panel: 저장 위치 & 파일 ──────────────────────────────
    panel(s, 0.62, 1.8, 6.05, 3.95, accent=BLUE)
    tf = textbox(s, 0.85, 1.96, 5.6, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [("저장 위치   ", 14, INK, True, FONT),
                  ("~/.agent-kanban/", 13, BLUE, True, MONO)], first=True)
    files = [
        ("active.json", "현재 보드 — 카드 전체"),
        ("archive/YYYY-MM.json", "월별 아카이브 (삭제 안 함)"),
        ("schedulers.json", "스케줄 (cron)"),
        ("scripts.json", "스크립트 (+ scripts/ 소스)"),
        ("settings.json", "API 키 · 설정"),
        ("telegram-state.json", "Telegram 채팅 상태"),
        ("screenshots/", "카드 첨부 이미지"),
    ]
    y = 2.55
    for name, desc in files:
        tn = textbox(s, 0.85, y, 2.85, 0.4, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        add_line(tn, [(name, 10.5, MONO_INK, False, MONO)], first=True)
        td = textbox(s, 3.72, y, 2.85, 0.4, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        add_line(td, [(desc, 10.5, MUT, False, FONT)], first=True)
        y += 0.44

    # ── right panel: 어떻게 저장되나 ──────────────────────────────
    panel(s, 6.88, 1.8, 5.83, 3.95, accent=GREEN)
    tf = textbox(s, 7.12, 1.96, 5.4, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [("어떻게 저장되나", 14, INK, True, FONT)], first=True)
    bullets = [
        [("도메인별 ", 11.5, INK, False, FONT), ("JSON 파일", 11.5, GREEN, True, FONT),
         (" — 별도 DB 없음", 11.5, INK, False, FONT)],
        [("원자적 쓰기: ", 11.5, INK, True, FONT),
         (".tmp 에 쓴 뒤 rename 으로 교체 → 쓰다 깨져도 원본 보존", 11.5, MUT, False, FONT)],
        [("듀얼 락: ", 11.5, INK, True, FONT),
         ("in-process mutex + 파일 락(.lock) → 플러그인·서버·CLI 동시 접근 직렬화", 11.5, MUT, False, FONT)],
        [("2-space 들여쓴 JSON", 11.5, INK, False, FONT),
         (" — 사람이 그대로 읽기 가능", 11.5, MUT, False, FONT)],
    ]
    tf = textbox(s, 7.12, 2.5, 5.45, 3.1)
    first = True
    for parts in bullets:
        add_line(tf, [("•  ", 11.5, GREEN, True, FONT)] + parts, first=first, space_after=9)
        first = False

    # ── stack panel ───────────────────────────────────────────────
    panel(s, 0.62, 5.95, 12.09, 1.05, accent=PURPLE)
    tf = textbox(s, 0.92, 5.95, 11.5, 1.05, anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [
        ("STACK    ", 13, PURPLE, True, FONT),
        ("Bun", 12, INK, True, FONT),
        (" (런타임·번들러·테스트러너)  ·  ", 12, MUT, False, FONT),
        ("TypeScript 5.9", 12, INK, True, FONT),
        ("  ·  ", 12, MUT, False, FONT),
        ("React 19 + react-dom", 12, INK, True, FONT),
        ("  ·  ", 12, MUT, False, FONT),
        ("Vite 7 (+plugin-react)", 12, INK, True, FONT),
        ("  ·  ", 12, MUT, False, FONT),
        ("@opencode-ai/plugin", 12, INK, True, FONT),
        ("  ·  ", 12, MUT, False, FONT),
        ("zod", 12, INK, True, FONT), (" 검증  ·  ", 12, MUT, False, FONT),
        ("nanoid", 12, INK, True, FONT), (" ID  ·  ", 12, MUT, False, FONT),
        ("croner", 12, INK, True, FONT), (" cron  ·  ", 12, MUT, False, FONT),
        ("@playwright/test", 12, INK, True, FONT), (" E2E", 12, MUT, False, FONT),
    ], first=True)


def main():
    prs = Presentation(TARGET)
    slides = list(prs.slides)
    s = slides[SLIDE_INDEX]
    txt = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
    if not ("코드" in txt or "저장" in txt):
        raise SystemExit(f"슬라이드 {SLIDE_INDEX + 1}이 예상과 다릅니다: {txt[:60]!r}")
    rebuild(s)
    prs.save(TARGET)
    print(f"saved: {TARGET}  (rebuilt slide {SLIDE_INDEX + 1} as storage + stack)")


if __name__ == "__main__":
    main()
